#!/usr/bin/env python3
"""Create an updated course deck with segmented fitting and counterfactual analysis."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


SOURCE = Path("presentations/digital_economy_numerical_20260519_113555.pptx")
OUTPUT = Path("presentations/digital_economy_numerical_segmented_20260521.pptx")


TEXT_REPLACEMENTS = {
    "全国博物馆数量增长": "数字经济背景下",
    "与文旅产业复苏": "博物馆增长与文旅复苏",
    "基于三点差分、最小二乘拟合与相关性分析": "文化供给、数字化连接与消费转化的数值分析",
    "选题价值：文旅融合是数字经济的民生侧入口": "选题价值：数字经济连接文化供给与文旅消费",
    "博物馆不只是文化机构数量指标，也反映公共服务供给、城市文旅吸引力和数字化展示能力的共同提升。": (
        "本文不只考察场馆数量，而是把博物馆看作文化资源入口，分析其在数字化连接下向游客流量和文旅消费转化的过程。"
    ),
    "本报告的问题：博物馆供给扩张与文旅需求恢复，是否呈现可量化的同向关系？": (
        "本报告的问题：数字经济背景下，文化供给扩张能否重新转化为游客恢复和文旅消费？"
    ),
    "数据预处理：统一尺度，并保留冲击信息": "数据处理：从原始指标到可计算特征",
    "社会经济数据不是光滑物理曲线，预处理的重点是减少口径误差，同时不抹平真实冲击。": (
        "原始数据只有 10 个年度点，因此处理重点不是扩充样本，而是统一口径、构造指数和转化效率指标。"
    ),
    "算法二：最小二乘拟合用于短期趋势估计": "算法二：分段最小二乘避免冲击期扭曲",
    "疫情冲击导致全周期拟合偏差较大，因此采用“平稳增长期解释 + 恢复期预测”的分段思路。": (
        "课程中的最小二乘用于最小化残差平方和；本报告比较全周期线性、二次与分段线性，"
        "用残差和留一误差说明分段低次模型更稳。"
    ),
    "模型选择比较": "模型选择与课程依据",
    "全周期模型": "全周期线性/二次",
    "优点：简单": "优点：公式简单",
    "缺点：把冲击和趋势": "缺点：把冲击和长期趋势",
    "混在一个斜率里": "混在同一模型里",
    "分段模型": "分段低次模型",
    "优点：解释机制更清楚": "符合课程分段低次思想",
    "恢复段斜率可用于": "能分别解释平稳增长",
    "短期预测": "和冲击后恢复",
    "本报告采用：2020-2023 年恢复段线性拟合，作为短期外推基础。": (
        "本报告采用：2014-2019 年作为平稳增长段，2020-2023 年作为冲击后恢复段；"
        "同时用平稳趋势构造反事实基准。"
    ),
    "误差意识：模型预测不是确定结论，而是基于历史恢复速度的情景估计。": (
        "误差意识：二次模型样本内略好但留一误差更差，因此不采用高阶外推。"
    ),
    "相关性：平稳期高度同向，冲击期削弱全周期关系": "模型比较：分段线性更适合冲击型社会数据",
    "Pearson 相关系数显示：2014-2019 年关系非常紧密，全周期因疫情冲击转弱甚至失真。": (
        "用 RMSE 和留一误差比较模型，说明分段低次最小二乘更稳。"
    ),
    "短期预测与误差分析：恢复段可外推，但不宜过度解释": "反事实缺口：衡量文旅需求相对正常趋势的恢复程度",
    "基于 2020-2023 年恢复段线性拟合，得到 2024-2026 年情景预测；结果用于趋势判断，不作精确预测。": (
        "以 2014-2019 年平稳增长趋势作为无冲击基准，对比 2020-2023 年真实游客与消费水平。"
    ),
    "恢复段拟合预测": "反事实缺口结果",
    "误差来源": "误差来源与边界",
    "误差分析结论：模型更适合解释方向和弹性，不适合替代统计部门正式预测。": (
        "结论：分段拟合用于解释结构变化，反事实缺口用于衡量恢复程度；两者都不替代正式统计预测。"
    ),
    "三点差分识别了 2020": "三点差分识别冲击，",
    "年冲击与 2023 年修复，": "分段最小二乘比较模型，",
    "分段拟合更合理。": "反事实缺口衡量恢复。",
    "结论：文化设施供给是文旅复苏的基础条件": "结论：数字经济提升文化供给的消费转化效率",
    "数字化能力提升，则决定供给向消费转化的效率": "博物馆数量增长提供供给基础，数字化连接决定转化效率",
    "公共文化设施扩张叠加": "线上预约、智慧导览与",
    "数字化展示，有助于释放": "数字展陈提升连接效率，",
    "文旅消费潜力。": "推动消费转化。",
    "汇报亮点：用数值方法解释“长期增长、短期冲击、恢复修复”三个阶段。": (
        "汇报亮点：用数值方法解释“文化供给扩张—数字化连接—文旅消费转化”的链条。"
    ),
}


def replace_existing_text(prs: Presentation) -> None:
    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            text = shape.text
            if text in TEXT_REPLACEMENTS:
                shape.text = TEXT_REPLACEMENTS[text]


def add_cover_box(slide, left, top, width, height) -> None:
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 250, 252)
    box.line.color.rgb = RGBColor(220, 226, 232)
    box.line.width = Pt(0.75)


def add_plain_cover(slide, left, top, width, height) -> None:
    box = slide.shapes.add_shape(1, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(247, 250, 252)
    box.line.fill.background()


def add_textbox(slide, left, top, width, height, text, font_size=15, bold=False, color=(32, 42, 54)) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    text_frame = box.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    run = paragraph.add_run()
    run.text = text
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)


def add_metric_card(slide, left, top, width, height, title, value, note) -> None:
    add_cover_box(slide, left, top, width, height)
    add_textbox(slide, left + Inches(0.14), top + Inches(0.10), width - Inches(0.28), Inches(0.25), title, 11, True, (75, 85, 99))
    add_textbox(slide, left + Inches(0.14), top + Inches(0.40), width - Inches(0.28), Inches(0.38), value, 22, True, (24, 99, 112))
    add_textbox(slide, left + Inches(0.14), top + Inches(0.86), width - Inches(0.28), Inches(0.34), note, 10, False, (75, 85, 99))


def add_chain_node(slide, left, top, width, title, note) -> None:
    add_cover_box(slide, left, top, width, Inches(0.62))
    add_textbox(slide, left + Inches(0.12), top + Inches(0.09), width - Inches(0.24), Inches(0.20), title, 10, True, (24, 99, 112))
    add_textbox(slide, left + Inches(0.12), top + Inches(0.33), width - Inches(0.24), Inches(0.20), note, 8, False, (75, 85, 99))


def update_slide_2(prs: Presentation) -> None:
    slide = prs.slides[1]
    add_plain_cover(slide, Inches(0.00), Inches(0.96), Inches(13.33), Inches(0.72))
    add_textbox(
        slide,
        Inches(0.75),
        Inches(1.20),
        Inches(11.90),
        Inches(0.28),
        "博物馆是文化资源入口，数字化连接决定其向游客流量和文旅消费的转化效率。",
        13,
        False,
        (32, 42, 54),
    )
    add_plain_cover(slide, Inches(1.15), Inches(5.08), Inches(11.00), Inches(1.25))
    add_textbox(slide, Inches(1.40), Inches(5.30), Inches(2.20), Inches(0.22), "数字经济机制", 10, True, (75, 85, 99))
    add_chain_node(slide, Inches(2.15), Inches(5.61), Inches(2.20), "供给层", "博物馆数量增长")
    add_textbox(slide, Inches(4.52), Inches(5.78), Inches(0.35), Inches(0.20), "→", 16, True, (24, 99, 112))
    add_chain_node(slide, Inches(4.95), Inches(5.61), Inches(2.70), "数字化连接层", "预约 / 导览 / 展陈 / 文创")
    add_textbox(slide, Inches(7.82), Inches(5.78), Inches(0.35), Inches(0.20), "→", 16, True, (24, 99, 112))
    add_chain_node(slide, Inches(8.25), Inches(5.61), Inches(2.35), "转化层", "游客恢复 + 消费支出")


def add_processing_step(slide, left, top, width, title, body) -> None:
    add_cover_box(slide, left, top, width, Inches(0.88))
    add_textbox(slide, left + Inches(0.14), top + Inches(0.11), width - Inches(0.28), Inches(0.22), title, 11, True, (24, 99, 112))
    add_textbox(slide, left + Inches(0.14), top + Inches(0.42), width - Inches(0.28), Inches(0.34), body, 9, False, (75, 85, 99))


def update_slide_5(prs: Presentation) -> None:
    slide = prs.slides[4]
    add_plain_cover(slide, Inches(0.00), Inches(0.96), Inches(13.33), Inches(0.72))
    add_textbox(
        slide,
        Inches(0.75),
        Inches(1.20),
        Inches(11.90),
        Inches(0.28),
        "原始数据虽少，但可通过口径统一、指数化和转化效率指标形成可计算时间序列。",
        13,
        False,
        (32, 42, 54),
    )
    add_cover_box(slide, Inches(0.55), Inches(1.55), Inches(12.20), Inches(4.85))
    add_textbox(slide, Inches(0.95), Inches(1.82), Inches(10.80), Inches(0.34), "处理链路：清洗整理不是扩大样本，而是形成适合计算的离散点集", 17, True)

    add_processing_step(slide, Inches(0.95), Inches(2.38), Inches(2.18), "1 统一口径", "年份对齐；统一单位")
    add_processing_step(slide, Inches(3.35), Inches(2.38), Inches(2.18), "2 指数化", "2014=100；同尺度比较")
    add_processing_step(slide, Inches(5.75), Inches(2.38), Inches(2.18), "3 转化指标", "单馆游客；人均消费")
    add_processing_step(slide, Inches(8.15), Inches(2.38), Inches(2.18), "4 阶段标注", "增长、冲击、修复、恢复")
    add_processing_step(slide, Inches(10.55), Inches(2.38), Inches(1.78), "5 模型输入", "差分 / 拟合 / 缺口")

    add_textbox(slide, Inches(0.95), Inches(3.70), Inches(3.60), Inches(0.28), "派生指标示例", 15, True)
    add_metric_card(slide, Inches(0.95), Inches(4.08), Inches(2.55), Inches(1.22), "单馆游客量", "1.170 → 0.716", "2019 到 2023，百万人次/馆")
    add_metric_card(slide, Inches(3.80), Inches(4.08), Inches(2.55), Inches(1.22), "人均旅游消费", "953.2 → 1004.6", "2019 到 2023，元/人")
    add_metric_card(slide, Inches(6.65), Inches(4.08), Inches(2.55), Inches(1.22), "博物馆指数", "100 → 186.8", "2014 到 2023，2014=100")
    add_textbox(
        slide,
        Inches(9.45),
        Inches(4.08),
        Inches(2.80),
        Inches(1.08),
        "供给扩张，但单馆游客未回 2019；\n人均消费更快恢复，体现转化效率。",
        9,
        False,
        (32, 42, 54),
    )
    add_textbox(slide, Inches(0.95), Inches(5.62), Inches(10.90), Inches(0.30), "输出：processed_features.csv + diff_results.csv + model_comparison.csv + counterfactual_gap.csv", 11, True, (24, 99, 112))


def update_slide_10(prs: Presentation) -> None:
    slide = prs.slides[9]
    add_cover_box(slide, Inches(0.00), Inches(1.55), Inches(13.33), Inches(5.35))
    add_textbox(slide, Inches(1.05), Inches(1.90), Inches(6.45), Inches(0.35), "模型误差比较（RMSE，越小越好）", 18, True)
    add_metric_card(slide, Inches(1.05), Inches(2.55), Inches(2.85), Inches(1.30), "国内游客", "432.8 vs 1082.5", "分段线性 RMSE / 全周期线性 RMSE")
    add_metric_card(slide, Inches(4.20), Inches(2.55), Inches(2.85), Inches(1.30), "旅游总花费", "5118.5 vs 11968.2", "分段线性 RMSE / 全周期线性 RMSE")
    add_metric_card(slide, Inches(7.35), Inches(2.55), Inches(2.85), Inches(1.30), "留一误差", "1747.7 > 1340.5", "游客二次模型 / 线性模型")
    add_textbox(
        slide,
        Inches(1.05),
        Inches(4.25),
        Inches(10.85),
        Inches(0.95),
        "二次模型虽然样本内 RMSE 略低，但国内游客和旅游消费的留一误差更高，说明小样本下高阶拟合不稳；因此采用分段低次最小二乘。",
        14,
    )
    add_textbox(slide, Inches(1.05), Inches(5.35), Inches(10.85), Inches(0.35), "课程对应：函数逼近 + 最小二乘残差平方和 + 分段低次思想", 13, True, (24, 99, 112))


def update_slide_11(prs: Presentation) -> None:
    slide = prs.slides[10]
    add_cover_box(slide, Inches(0.80), Inches(1.80), Inches(6.25), Inches(3.95))
    add_textbox(slide, Inches(1.05), Inches(2.05), Inches(5.80), Inches(0.35), "实际值 / 2014-2019 平稳趋势", 18, True)
    add_metric_card(slide, Inches(1.05), Inches(2.65), Inches(2.65), Inches(1.28), "国内游客 2020", "44.4%", "缺口 -3602.5 百万人次")
    add_metric_card(slide, Inches(4.00), Inches(2.65), Inches(2.65), Inches(1.28), "国内游客 2023", "61.5%", "仍低于平稳趋势")
    add_metric_card(slide, Inches(1.05), Inches(4.20), Inches(2.65), Inches(1.28), "旅游消费 2020", "35.8%", "缺口 -39949.8 亿元")
    add_metric_card(slide, Inches(4.00), Inches(4.20), Inches(2.65), Inches(1.28), "旅游消费 2023", "62.4%", "恢复明显但未回趋势线")

    replacements = {
        "年度采样间隔较大": "年度采样间隔较大",
        "冲击事件使序列非平稳": "平稳趋势本身也有模型误差",
        "博物馆数与旅游需求存在滞后效应": "外生冲击和政策变化会改变恢复速度",
    }
    for shape in slide.shapes:
        if hasattr(shape, "text") and shape.text in replacements:
            shape.text = replacements[shape.text]

    add_plain_cover(slide, Inches(1.05), Inches(6.05), Inches(10.95), Inches(0.50))
    add_textbox(
        slide,
        Inches(1.18),
        Inches(6.15),
        Inches(10.60),
        Inches(0.32),
        "结论：分段拟合用于解释结构变化，反事实缺口用于衡量恢复程度；两者都不替代正式统计预测。",
        14,
        False,
        (32, 42, 54),
    )


def main() -> None:
    prs = Presentation(SOURCE)
    replace_existing_text(prs)
    update_slide_2(prs)
    update_slide_5(prs)
    update_slide_10(prs)
    update_slide_11(prs)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    main()
